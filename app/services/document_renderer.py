"""新建办公文档的确定性渲染器。

模型只允许描述内容，不能执行代码或选择文件系统路径。本模块校验精简的文档
规格，并将可复核产物写入调用方提供的输出目录。
"""

from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import Any


SUPPORTED_FORMATS = {"docx", "pptx", "xlsx"}
SUPPORTED_STYLES = {"business", "minimal", "academic", "modern"}
MAX_SECTIONS = 30
MAX_SLIDES = 30
MAX_SHEETS = 10
MAX_ROWS = 500
MAX_COLUMNS = 30


_STYLE_PALETTES = {
    "business": {"primary": "1F4E78", "accent": "2F75B5", "text": "1F2937", "light": "EAF2F8"},
    "minimal": {"primary": "374151", "accent": "6B7280", "text": "111827", "light": "F3F4F6"},
    "academic": {"primary": "7C2D12", "accent": "B45309", "text": "292524", "light": "FFF7ED"},
    "modern": {"primary": "0F766E", "accent": "14B8A6", "text": "134E4A", "light": "F0FDFA"},
}


def _text(value: object, limit: int = 4000) -> str:
    """Normalize untrusted model text before inserting it into a document."""
    text = str(value or "").replace("\x00", "").strip()
    return text[:limit]


def _text_list(value: object, limit: int = 40, item_limit: int = 2000) -> list[str]:
    if not isinstance(value, list):
        return []
    return [item for item in (_text(item, item_limit) for item in value[:limit]) if item]


def safe_output_filename(value: object, document_format: str) -> str:
    """Return a basename with the requested extension, never a server path."""
    raw = str(value or "").strip().replace("\\", "/")
    if not raw or "/" in raw:
        raise ValueError("文件名必须是单个文件名，不能包含路径")
    name = Path(raw).name
    if name in {".", ".."} or not re.fullmatch(r"[\w.\-\u4e00-\u9fff]{1,120}", name):
        raise ValueError("文件名仅可包含中文、字母、数字、下划线、短横线和点")
    suffix = f".{document_format}"
    if Path(name).suffix.casefold() != suffix:
        name = f"{Path(name).stem or 'document'}{suffix}"
    return name


def normalize_spec(params: dict[str, Any]) -> dict[str, Any]:
    """Validate the finite renderer input contract.

    Unknown fields are deliberately ignored.  They are reserved for future
    template support and must not become an unreviewed rendering surface.
    """
    raw_format = str(params.get("format") or "").casefold().lstrip(".")
    if raw_format not in SUPPORTED_FORMATS:
        raise ValueError("format 仅支持 docx、pptx 或 xlsx")
    style = str(params.get("style") or "business").casefold()
    if style not in SUPPORTED_STYLES:
        style = "business"
    title = _text(params.get("title"), 200) or "未命名文档"
    return {
        "format": raw_format,
        "filename": safe_output_filename(params.get("filename"), raw_format),
        "title": title,
        "style": style,
        "sections": _normalize_sections(params.get("sections")),
        "slides": _normalize_slides(params.get("slides")),
        "sheets": _normalize_sheets(params.get("sheets")),
    }


def _normalize_table(value: object) -> dict[str, list] | None:
    if not isinstance(value, dict):
        return None
    headers = _text_list(value.get("headers"), MAX_COLUMNS, 200)
    rows: list[list[str]] = []
    raw_rows = value.get("rows")
    if isinstance(raw_rows, list):
        for row in raw_rows[:MAX_ROWS]:
            if not isinstance(row, list):
                continue
            rows.append([_text(cell, 500) for cell in row[:MAX_COLUMNS]])
    if not headers and not rows:
        return None
    width = len(headers) or max((len(row) for row in rows), default=1)
    width = min(max(width, 1), MAX_COLUMNS)
    headers = (headers + [""] * width)[:width]
    rows = [(row + [""] * width)[:width] for row in rows]
    return {"headers": headers, "rows": rows}


def _normalize_sections(value: object) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    sections = []
    for item in value[:MAX_SECTIONS]:
        if not isinstance(item, dict):
            continue
        heading = _text(item.get("heading"), 240)
        paragraphs = _text_list(item.get("paragraphs"), 25)
        bullets = _text_list(item.get("bullets"), 30, 1000)
        table = _normalize_table(item.get("table"))
        if heading or paragraphs or bullets or table:
            sections.append({"heading": heading, "paragraphs": paragraphs, "bullets": bullets, "table": table})
    return sections


def _normalize_slides(value: object) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    slides = []
    for item in value[:MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        title = _text(item.get("title"), 160)
        subtitle = _text(item.get("subtitle"), 500)
        bullets = _text_list(item.get("bullets"), 8, 360)
        table = _normalize_table(item.get("table"))
        if title or subtitle or bullets or table:
            slides.append({"title": title, "subtitle": subtitle, "bullets": bullets, "table": table})
    return slides


def _normalize_sheets(value: object) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    sheets = []
    for index, item in enumerate(value[:MAX_SHEETS], start=1):
        if not isinstance(item, dict):
            continue
        name = _text(item.get("name"), 31) or f"Sheet{index}"
        name = re.sub(r"[\\/*?:\[\]]", "_", name)[:31]
        table = _normalize_table({"headers": item.get("headers"), "rows": item.get("rows")})
        if table:
            sheets.append({"name": name, **table})
    return sheets


def render_document(params: dict[str, Any], output_dir: Path) -> Path:
    """Render a validated office document into one user-isolated directory."""
    spec = normalize_spec(params)
    destination_dir = Path(output_dir).resolve()
    destination_dir.mkdir(parents=True, exist_ok=True)
    target = (destination_dir / spec["filename"]).resolve()
    if target.parent != destination_dir:
        raise ValueError("文件输出目录无效")
    with tempfile.NamedTemporaryFile(prefix=".lumi-", suffix=target.suffix, dir=destination_dir, delete=False) as handle:
        temporary = Path(handle.name)
    try:
        if spec["format"] == "docx":
            _render_docx(spec, temporary)
        elif spec["format"] == "pptx":
            _render_pptx(spec, temporary)
        else:
            _render_xlsx(spec, temporary)
        if not temporary.is_file() or temporary.stat().st_size == 0:
            raise ValueError("文档渲染未生成有效文件")
        temporary.replace(target)
        return target
    except Exception:
        temporary.unlink(missing_ok=True)
        raise


def _render_docx(spec: dict[str, Any], target: Path) -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    palette = _STYLE_PALETTES[spec["style"]]
    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
    normal = doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal.font.size = Pt(10.5)
    title = doc.add_heading(spec["title"], level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.name = "Microsoft YaHei"
        run.font.color.rgb = RGBColor.from_string(palette["primary"])
    sections = spec["sections"] or [{"heading": "内容", "paragraphs": ["根据你的要求生成。"], "bullets": [], "table": None}]
    for section in sections:
        if section["heading"]:
            doc.add_heading(section["heading"], level=1)
        for paragraph in section["paragraphs"]:
            doc.add_paragraph(paragraph)
        for bullet in section["bullets"]:
            doc.add_paragraph(bullet, style="List Bullet")
        if section["table"]:
            _append_docx_table(doc, section["table"], palette["primary"])
    doc.save(str(target))


def _append_docx_table(doc, table_spec: dict[str, list], primary: str) -> None:
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import RGBColor

    headers = table_spec["headers"]
    rows = table_spec["rows"]
    table = doc.add_table(rows=1 if headers else 0, cols=len(headers) or len(rows[0]))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    if headers:
        for index, value in enumerate(headers):
            cell = table.rows[0].cells[index]
            cell.text = value
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            shade = OxmlElement("w:shd")
            shade.set(qn("w:fill"), primary)
            cell._tc.get_or_add_tcPr().append(shade)
            for run in cell.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
    for values in rows:
        cells = table.add_row().cells
        for index, value in enumerate(values):
            cells[index].text = value


def _render_pptx(spec: dict[str, Any], target: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    palette = _STYLE_PALETTES[spec["style"]]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides = spec["slides"] or [{"title": spec["title"], "subtitle": "根据你的要求生成", "bullets": [], "table": None}]
    for index, slide_spec in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _paint_slide(slide, palette["light"])
        if index == 0 and not slide_spec["bullets"] and not slide_spec["table"]:
            _add_ppt_text(slide, slide_spec["title"] or spec["title"], 0.9, 2.2, 11.5, 0.8, 32, palette["primary"], bold=True)
            _add_ppt_text(slide, slide_spec["subtitle"], 0.95, 3.2, 11.2, 1.0, 18, palette["text"])
            continue
        _add_ppt_text(slide, slide_spec["title"] or spec["title"], 0.75, 0.55, 11.9, 0.6, 25, palette["primary"], bold=True)
        if slide_spec["bullets"]:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(6.3 if slide_spec["table"] else 11.0), Inches(4.9))
            frame = box.text_frame
            frame.clear()
            for bullet_index, bullet in enumerate(slide_spec["bullets"]):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                paragraph.text = bullet
                paragraph.level = 0
                paragraph.font.size = Pt(19)
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.color.rgb = RGBColor.from_string(palette["text"])
                paragraph.space_after = Pt(12)
        elif slide_spec["subtitle"]:
            _add_ppt_text(slide, slide_spec["subtitle"], 1.0, 1.55, 10.9, 3.8, 20, palette["text"])
        if slide_spec["table"]:
            _append_ppt_table(slide, slide_spec["table"], 7.35, 1.55, 5.2, 4.6, palette)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.75), Inches(2.1), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string(palette["accent"])
        accent.line.fill.background()
    prs.save(str(target))


def _paint_slide(slide, color: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _add_ppt_text(slide, value: str, left: float, top: float, width: float, height: float, size: int, color: str, bold: bool = False) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = value
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.name = "Microsoft YaHei"
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)


def _append_ppt_table(slide, table_spec: dict[str, list], left: float, top: float, width: float, height: float, palette: dict[str, str]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    headers, rows = table_spec["headers"], table_spec["rows"]
    row_count = max(1, len(rows) + (1 if headers else 0))
    column_count = max(1, len(headers) or len(rows[0]))
    shape = slide.shapes.add_table(row_count, column_count, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    values = ([headers] if headers else []) + rows
    for row_index, row in enumerate(values):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(palette["primary"] if row_index == 0 and headers else "FFFFFF")
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.name = "Microsoft YaHei"
                    run.font.color.rgb = RGBColor(255, 255, 255) if row_index == 0 and headers else RGBColor.from_string(palette["text"])


def _render_xlsx(spec: dict[str, Any], target: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    palette = _STYLE_PALETTES[spec["style"]]
    workbook = Workbook()
    first = workbook.active
    sheets = spec["sheets"] or [{"name": "Sheet1", "headers": [spec["title"]], "rows": []}]
    for index, sheet_spec in enumerate(sheets):
        sheet = first if index == 0 else workbook.create_sheet()
        sheet.title = sheet_spec["name"]
        headers, rows = sheet_spec["headers"], sheet_spec["rows"]
        if headers:
            sheet.append(headers)
            for cell in sheet[1]:
                cell.font = Font(name="Microsoft YaHei", bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor=palette["primary"])
                cell.alignment = Alignment(horizontal="center", vertical="center")
        for row in rows:
            sheet.append(row)
        sheet.freeze_panes = "A2" if headers else "A1"
        for column_index in range(1, max(len(headers), max((len(row) for row in rows), default=0)) + 1):
            values = [str(sheet.cell(row=row_index, column=column_index).value or "") for row_index in range(1, min(sheet.max_row, 80) + 1)]
            sheet.column_dimensions[get_column_letter(column_index)].width = min(max(max((len(v) for v in values), default=8) + 2, 10), 36)
    workbook.save(str(target))
