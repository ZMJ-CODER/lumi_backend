"""办公文档结构化编辑引擎.

能力：
  - docx / xlsx / pptx：结构化读取与编辑（python-docx / openpyxl / python-pptx），
    不触碰二进制流，按"结构→编辑指令→应用"工作；
  - md / txt / json / csv / yaml 等纯文本：全量重写或 SEARCH/REPLACE 补丁；
  - 缓冲会话：编辑只作用于 buffered 副本，用户审核通过后取回落盘，可随时丢弃。

会话布局：data/office/{user_id}/{doc_id}/original.ext + buffered.ext + meta.json
"""

import csv
import json
import hashlib
import shutil
import time
import uuid
from datetime import datetime, timedelta, timezone
from pathlib import Path

from loguru import logger

from app.core.config import settings

OFFICE_DIR = Path(settings.UPLOAD_DIR).parent / "office"


# ── 通用脚本产物（新建文件，无源文档）─────────────────

def _safe_conv(conv_id: str) -> str:
    return "".join(ch for ch in str(conv_id or "default") if ch.isalnum() or ch in "-_")[:64] or "default"


def generic_outputs_dir(user_id: str, conv_id: str) -> Path:
    """通用产物目录：data/uploads/office_outputs/{user_id}/{conv_id}."""
    return Path(settings.UPLOAD_DIR) / "office_outputs" / str(user_id) / _safe_conv(conv_id)


def list_generic_outputs(user_id: str, conv_id: str) -> list[dict]:
    """列出某任务/会话的通用产物."""
    d = generic_outputs_dir(user_id, conv_id)
    if not d.exists():
        return []
    return [
        {"name": f.name, "size": f.stat().st_size}
        for f in sorted(d.iterdir())
        if f.is_file()
    ]


def resolve_generic_output(user_id: str, conv_id: str, name: str) -> Path | None:
    """解析通用产物路径（路径越权校验）."""
    d = generic_outputs_dir(user_id, conv_id).resolve()
    try:
        target = (d / str(name)).resolve()
    except (ValueError, OSError):
        return None
    if d not in target.parents or not target.is_file():
        return None
    return target


def delete_generic_output(user_id: str, conv_id: str, name: str) -> bool:
    """删除通用产物（客户端投递到下载目录成功后调用）."""
    target = resolve_generic_output(user_id, conv_id, name)
    if target is None:
        return False
    try:
        target.unlink(missing_ok=True)
        return True
    except OSError:
        return False


def preview_generated_output(path: Path) -> dict:
    """生成可安全内嵌展示的产物摘要，绝不返回磁盘路径或原始二进制。"""
    suffix = path.suffix.lower()
    result: dict = {"name": path.name, "size": path.stat().st_size, "kind": suffix.lstrip(".")}
    max_chars = 60_000
    max_rows, max_cols = 50, 30
    if suffix in {".csv", ".tsv"}:
        delimiter = "\t" if suffix == ".tsv" else ","
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as fh:
            rows = []
            for index, row in enumerate(csv.reader(fh, delimiter=delimiter)):
                if index >= max_rows:
                    result["truncated"] = True
                    break
                rows.append([str(cell)[:200] for cell in row[:max_cols]])
        return {**result, "preview_type": "table", "rows": rows}
    if suffix in {".xlsx", ".xlsm"}:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(path, read_only=True, data_only=False)
            sheet = wb[wb.sheetnames[0]] if wb.sheetnames else None
            rows = []
            if sheet is not None:
                for index, row in enumerate(sheet.iter_rows(values_only=True)):
                    if index >= max_rows:
                        result["truncated"] = True
                        break
                    rows.append(["" if cell is None else str(cell)[:200] for cell in row[:max_cols]])
            result["sheets"] = wb.sheetnames[:20]
            return {**result, "preview_type": "table", "rows": rows}
        except Exception:
            return {**result, "preview_type": "unavailable", "message": "此表格暂无法预览，请下载后查看。"}
    if suffix == ".docx":
        try:
            from docx import Document

            document = Document(str(path))
            lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
            for table in document.tables[:8]:
                for row in table.rows[:20]:
                    lines.append(" | ".join(cell.text.strip() for cell in row.cells))
            text = "\n".join(lines)
            return {**result, "preview_type": "text", "text": text[:max_chars], "truncated": len(text) > max_chars}
        except Exception:
            return {**result, "preview_type": "unavailable", "message": "此 Word 文件暂无法预览，请下载后查看。"}
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            presentation = Presentation(str(path))
            lines = []
            for index, slide in enumerate(presentation.slides[:30], start=1):
                values = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
                if values:
                    lines.append(f"第 {index} 页\n" + "\n".join(values))
            text = "\n\n".join(lines)
            return {**result, "preview_type": "text", "text": text[:max_chars], "truncated": len(text) > max_chars}
        except Exception:
            return {**result, "preview_type": "unavailable", "message": "此 PowerPoint 文件暂无法预览，请下载后查看。"}
    if suffix in {".txt", ".md", ".log", ".json", ".xml", ".yaml", ".yml", ".html", ".htm"}:
        with path.open("r", encoding="utf-8", errors="replace") as fh:
            text = fh.read(max_chars + 1)
        return {
            **result,
            "preview_type": "text",
            "text": text[:max_chars],
            "truncated": len(text) > max_chars,
        }
    return {**result, "preview_type": "unavailable", "message": "该文件格式不提供内嵌预览，请下载后查看。"}


def cleanup_generic_outputs(ttl_days: int = 7) -> int:
    """定时清理超过 TTL 的通用脚本产物目录（按用户 × 任务隔离）."""
    base = Path(settings.UPLOAD_DIR) / "office_outputs"
    if not base.exists():
        return 0
    cutoff = time.time() - max(int(ttl_days or 7), 1) * 86400
    removed = 0
    for user_dir in base.iterdir():
        if not user_dir.is_dir():
            continue
        for conv_dir in user_dir.iterdir():
            if not conv_dir.is_dir():
                continue
            try:
                if conv_dir.stat().st_mtime < cutoff:
                    shutil.rmtree(conv_dir, ignore_errors=True)
                    removed += 1
            except OSError:
                continue
    return removed
TEXT_EXTS = {".md", ".txt", ".json", ".csv", ".yaml", ".yml", ".toml", ".ini", ".log", ".xml"}
STRUCTURED_EXTS = {".docx", ".xlsx", ".pptx", ".docm", ".xlsm", ".pptm"}
# 老版 Office / 富文本：走本机 Office COM 提取（Windows + 已安装 Word/Excel/PowerPoint）
LEGACY_EXTS = {".doc", ".xls", ".ppt", ".rtf"}
# 需 Docling 解析的格式（PDF / ODT 等）：可读可分析，暂不支持结构化编辑
DOCLING_EXTS = {".pdf", ".odt"}
# 图片（扫描件/发票/截图）：Docling OCR 提取文本，可读可分析，不支持结构化编辑
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}
MAX_OFFICE_SIZE = 30 * 1024 * 1024

# 各格式允许的编辑操作（LLM 规划指令时用）
OP_SCHEMAS = {
    "docx": [
        '{"op":"replace_paragraph","index":0,"text":"新内容"}',
        '{"op":"add_paragraph","text":"新段落","index":null}',
        '{"op":"delete_paragraph","index":0}',
        '{"op":"replace_all","find":"旧词","replace":"新词"}',
        '{"op":"replace_table_cell","table_index":0,"row":0,"col":0,"text":"新内容"}',
        '{"op":"add_table_row","table_index":0,"values":["a","b"]}',
    ],
    "xlsx": [
        '{"op":"set_cell","sheet":"Sheet1","cell":"A1","value":"x"}',
        '{"op":"replace_all","sheet":"Sheet1","find":"旧值","replace":"新值"}',
        '{"op":"add_row","sheet":"Sheet1","values":["a","b"]}',
    ],
    "pptx": [
        '{"op":"replace_text","slide_index":0,"find":"旧词","replace":"新词"}',
        '{"op":"add_text","slide_index":0,"text":"新文本"}',
        '{"op":"delete_slide","slide_index":0}',
        '{"op":"add_slide","title":"新幻灯片标题","text":"正文内容"}',
    ],
    "text": [
        '{"op":"rewrite","content":"完整新内容"}',
        '{"op":"search_replace","old":"原文片段","new":"新片段","first_only":false}',
    ],
}


async def plan_edits(
    instruction: str,
    structure: str,
    kind: str,
    user_id: str,
    api_key: str | None = None,
) -> list[dict]:
    """让 LLM 根据文档结构与用户指令，输出结构化的编辑指令 JSON 列表."""
    from app.core.llm import LLMClient
    from app.services.usage import CATEGORY_SKILL

    llm = LLMClient()
    ops = OP_SCHEMAS.get(kind, OP_SCHEMAS["text"])
    system = (
        "你是办公文档编辑规划器。根据文档结构和用户指令，输出要执行的编辑操作列表。\n"
        f"本文件类型 {kind} 允许的操作（字段按示例）：\n"
        + "\n".join(f"- {o}" for o in ops)
        + "\n约束：索引从 0 开始；只输出 JSON 数组（不要 Markdown 围栏、不要解释）；"
        "无法满足的指令输出空数组 []。"
    )
    from app.core.agent_security import UNTRUSTED_CONTENT_RULES

    system += "\n\n" + UNTRUSTED_CONTENT_RULES
    reply = await llm.chat(
        [
            {"role": "system", "content": system},
            {"role": "user", "content": f"用户指令：{instruction}\n\n文档结构：\n{structure[:60000]}"},
        ],
        scene="office",
        max_tokens=4000,
        temperature=0.1,
        usage_user_id=user_id,
        usage_category=CATEGORY_SKILL,
        disable_reasoning_effort=True,
        api_key=api_key,
    )
    text = (reply or "").strip()
    if text.startswith("```"):
        text = text.strip("`")
        if text.startswith("json"):
            text = text[4:]
    start, end = text.find("["), text.rfind("]")
    if start < 0 or end <= start:
        return []
    try:
        data = json.loads(text[start : end + 1])
        return data if isinstance(data, list) else []
    except (ValueError, TypeError):
        return []


def _session_dir(user_id: str, doc_id: str) -> Path:
    safe_u = "".join(c for c in str(user_id) if c.isalnum() or c in "-_") or "anon"
    safe_d = "".join(c for c in str(doc_id) if c.isalnum() or c in "-_")
    return OFFICE_DIR / safe_u / safe_d


def _original_path(session: Path, ext: str) -> Path:
    return session / f"original{ext}"


def _buffered_path(session: Path, ext: str) -> Path:
    return session / f"buffered{ext}"


def _meta_path(session: Path) -> Path:
    return session / "meta.json"


def detect_kind(filename: str) -> str:
    ext = Path(filename or "").suffix.lower()
    if ext in STRUCTURED_EXTS:
        return ext.lstrip(".")
    if ext in LEGACY_EXTS:
        return ext.lstrip(".")
    if ext in DOCLING_EXTS:
        return ext.lstrip(".")
    if ext in IMAGE_EXTS:
        return ext.lstrip(".")
    return "text"


def _legacy_extract_text(path: Path, kind: str) -> str:
    """老版 Office / 富文本提取全文.

    doc/rtf：纯 Python（OLE WordDocument 分段表 / RTF 剥离），跨环境可靠；
    xls/ppt：本机 Office COM（Windows，需安装 Office），失败给出明确提示。
    """
    if kind == "doc":
        text = _extract_doc_text(path)
        if text.strip():
            return text.strip()
        # 纯解析无结果（加密/损坏）：兜底尝试本机 Word COM
        return _extract_legacy_com(path, "doc")
    if kind == "rtf":
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(path.read_text(encoding="latin-1", errors="replace")) or ""
    return _extract_legacy_com(path, kind)


def _extract_doc_text(path: Path) -> str:
    """纯 Python 提取 Word 97-2003 (.doc) 文本：解析 OLE WordDocument 流的分段表（piece table）."""
    raw = path.read_bytes()
    if not raw:
        return ""
    # 伪装成 .doc 的 RTF / 纯文本
    if raw[:5].lower().startswith(b"{\\rtf"):
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(raw.decode("latin-1", errors="replace")) or ""
    try:
        text = raw.decode("utf-8")
        if text and (text.isprintable() or "\n" in text or "\r" in text):
            return text
    except Exception:  # noqa: BLE001
        pass
    # 非 OLE：尝试常见中文编码
    if not raw.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        for enc in ("gb18030", "utf-8-sig", "big5"):
            try:
                return raw.decode(enc)
            except Exception:  # noqa: BLE001
                continue
        return ""
    try:
        import olefile
        import struct

        if not olefile.isOleFile(str(path)):
            return ""
        ole = olefile.OleFileIO(str(path))
        try:
            if not ole.exists("WordDocument"):
                return ""
            wd = ole.openstream("WordDocument").read()
            if len(wd) < 0x01AA:
                return ""
            fc_clx = struct.unpack("<I", wd[0x01A2:0x01A6])[0]
            lcb_clx = struct.unpack("<I", wd[0x01A6:0x01AA])[0]
            # Word 97+：分段表（CLX）在 0Table/1Table 流中，偏移相对该流
            which_table = struct.unpack("<H", wd[0x0A:0x0C])[0] & 0x0200
            table_name = "1Table" if which_table else "0Table"
            if not ole.exists(table_name):
                return ""
            table = ole.openstream(table_name).read()
            if fc_clx + lcb_clx > len(table):
                return ""
            clx = table[fc_clx : fc_clx + lcb_clx]
            plc_pcd = None
            pos = 0
            while pos < len(clx):
                t = clx[pos]
                if t == 0x02:  # Pcdt（分段表）
                    lcb = struct.unpack("<I", clx[pos + 1 : pos + 5])[0]
                    plc_pcd = clx[pos + 5 : pos + 5 + lcb]
                    break
                if t == 0x01:  # RgFc
                    cb = struct.unpack("<I", clx[pos + 1 : pos + 5])[0]
                    pos += 5 + cb
                else:
                    break
            if not plc_pcd or len(plc_pcd) < 4:
                return ""
            n = (len(plc_pcd) - 4) // 12
            if n <= 0:
                return ""
            cps = [struct.unpack("<I", plc_pcd[i * 4 : i * 4 + 4])[0] for i in range(n + 1)]
            parts = []
            for i in range(n):
                start, end = cps[i], cps[i + 1]
                if end <= start:
                    continue
                pcd = plc_pcd[(n + 1) * 4 + i * 8 : (n + 1) * 4 + (i + 1) * 8]
                flags = struct.unpack("<H", pcd[:2])[0]
                fc = struct.unpack("<I", pcd[2:6])[0] & 0x3FFFFFFF
                if fc >= len(wd):
                    continue
                if flags & 0x0004:  # 压缩（单字节）
                    fc //= 2
                    chunk = wd[fc : fc + (end - start)]
                    parts.append(chunk.decode("latin-1", errors="replace"))
                else:  # Unicode（UTF-16LE）
                    chunk = wd[fc : fc + (end - start) * 2]
                    parts.append(chunk.decode("utf-16-le", errors="replace"))
            return "".join(parts).strip()
        finally:
            ole.close()
    except Exception:  # noqa: BLE001
        return ""


def _extract_legacy_com(path: Path, kind: str) -> str:
    """老版 Office 用本机 Office COM 提取全文（Windows，需安装 Office）."""
    try:
        import pythoncom
        from win32com.client import gencache
    except ImportError as exc:  # noqa: BLE001
        raise ValueError(
            f"{kind.upper()} 格式解析需要本机安装 Office（win32com 不可用），请另存为 "
            + (".docx" if kind == "doc" else ".xlsx" if kind == "xls" else ".pptx" if kind == "ppt" else ".txt")
        ) from exc

    pythoncom.CoInitialize()
    try:
        if kind == "xls":
            app = gencache.EnsureDispatch("Excel.Application")
            app.Visible = False
            app.DisplayAlerts = False
            try:
                wb = app.Workbooks.Open(str(path), ReadOnly=True)
                try:
                    parts = []
                    for ws in wb.Worksheets:
                        parts.append(f"[Sheet: {ws.Name}]")
                        used = ws.UsedRange
                        for row in used.Rows:
                            vals = []
                            for cell in row.Cells:
                                v = cell.Value
                                vals.append("" if v is None else str(v))
                            if any(vals):
                                parts.append(" | ".join(vals))
                    return "\n".join(parts).strip()
                finally:
                    wb.Close(False)
            finally:
                app.Quit()
        if kind == "ppt":
            app = gencache.EnsureDispatch("PowerPoint.Application")
            try:
                prs = app.Presentations.Open(str(path), ReadOnly=True, WithWindow=False)
                try:
                    parts = []
                    for si, slide in enumerate(prs.Slides):
                        parts.append(f"[第{si}页]")
                        for shape in slide.Shapes:
                            try:
                                if shape.HasTextFrame and shape.TextFrame.HasText:
                                    parts.append(shape.TextFrame.TextRange.Text.strip())
                            except Exception:  # noqa: BLE001
                                continue
                    return "\n".join(parts).strip()
                finally:
                    prs.Close()
            finally:
                app.Quit()
    except Exception as exc:  # noqa: BLE001
        raise ValueError(
            f"{kind.upper()} 格式解析失败（{str(exc)[:120]}），请另存为新版格式后重试"
        ) from exc
    finally:
        try:
            pythoncom.CoUninitialize()
        except Exception:  # noqa: BLE001
            pass
    raise ValueError(f"不支持的旧格式: {kind}")


# ── 会话管理 ─────────────────────────────────────────────

def create_session(user_id: str, filename: str, content: bytes) -> dict:
    """保存上传文件并建立编辑会话；返回会话元数据."""
    if len(content) > MAX_OFFICE_SIZE:
        raise ValueError("文件超过 30MB 限制")
    ext = Path(filename).suffix.lower() or ".txt"
    doc_id = uuid.uuid4().hex[:12]
    session = _session_dir(user_id, doc_id)
    session.mkdir(parents=True, exist_ok=True)
    (session / f"original{ext}").write_bytes(content)
    meta = {
        "doc_id": doc_id,
        "filename": filename,
        "ext": ext,
        "kind": detect_kind(filename),
        "created_at": None,
        "committed": False,
    }
    _meta_path(session).write_text(json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8")
    return meta


def load_session(user_id: str, doc_id: str) -> dict:
    session = _session_dir(user_id, doc_id)
    if not session.exists():
        raise LookupError("办公文档会话不存在")
    meta = json.loads(_meta_path(session).read_text(encoding="utf-8"))
    meta["_session"] = str(session)
    return meta


# ── 临时会话 DB 持久化（聊天框上传链路） ────────────────
# 共享 Postgres：Docker / 本地后端切换时磁盘会话目录可能不在，但从 DB 可重建。

def _session_expire_at() -> datetime:
    return datetime.now(timezone.utc) + timedelta(hours=settings.OFFICE_SESSION_TTL_HOURS)


async def _refresh_session_use(user_id: str, doc_id: str) -> None:
    """会话被使用：刷新 last_used_at 与 expire_at（失败不阻塞主流程）."""
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy import update

        async with async_session_factory() as db:
            await db.execute(
                update(OfficeSession)
                .where(
                    OfficeSession.doc_id == str(doc_id),
                    OfficeSession.user_id == uuid.UUID(str(user_id)),
                )
                .values(last_used_at=datetime.now(timezone.utc), expire_at=_session_expire_at())
            )
            await db.commit()
    except Exception:  # noqa: BLE001
        pass


async def persist_session_record(
    user_id: str, meta: dict, content: bytes, content_text: str | None = None
) -> None:
    """上传后把临时会话写入 DB（提取全文 + 原始文件），供跨实例恢复."""
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy.dialects.postgresql import insert as pg_insert

        text = content_text if content_text is not None else extract_full_text(user_id, meta["doc_id"])
        async with async_session_factory() as db:
            stmt = pg_insert(OfficeSession).values(
                doc_id=meta["doc_id"],
                user_id=uuid.UUID(str(user_id)),
                filename=meta["filename"],
                kind=meta["kind"],
                content_text=text or None,
                file_bytes=content,
                file_hash=hashlib.sha256(content).hexdigest(),
                status="active",
                expire_at=_session_expire_at(),
                last_used_at=datetime.now(timezone.utc),
            )
            stmt = stmt.on_conflict_do_update(
                index_elements=[OfficeSession.doc_id],
                set_={
                    "filename": stmt.excluded.filename,
                    "kind": stmt.excluded.kind,
                    "content_text": stmt.excluded.content_text,
                    "file_bytes": stmt.excluded.file_bytes,
                    "file_hash": stmt.excluded.file_hash,
                    "status": "active",
                    "expire_at": stmt.excluded.expire_at,
                    "last_used_at": stmt.excluded.last_used_at,
                },
            )
            await db.execute(stmt)
            await db.commit()
    except Exception as exc:  # noqa: BLE001
        logger.warning("办公会话 DB 持久化失败: {}", exc)


async def promote_session_to_knowledge(
    user_id: str, doc_id: str, space_id: str, category: str | None = None
) -> dict:
    """在用户明确操作后将临时附件复制到长期知识库。

    该操作不自动触发，也不删除办公会话；同空间同 SHA-256 文件沿用知识库
    的幂等去重规则。正文和原始字节均来自会话本身，不能接受客户端路径。
    """
    from app.core.database import async_session_factory
    from app.models.db_models import OfficeSession
    from app.services.rag.knowledge import upload_document_file
    from sqlalchemy import select

    uid = uuid.UUID(str(user_id))
    async with async_session_factory() as db:
        row = (
            await db.execute(
                select(OfficeSession).where(
                    OfficeSession.doc_id == str(doc_id),
                    OfficeSession.user_id == uid,
                    OfficeSession.status == "active",
                )
            )
        ).scalar_one_or_none()
        if row is None:
            raise LookupError("办公文档会话不存在")
        content = row.file_bytes
        if not content:
            meta = await ensure_session(user_id, doc_id)
            content = (Path(meta["_session"]) / f"original{meta['ext']}").read_bytes()
        doc, path, is_new = await upload_document_file(
            db, user_id, space_id, row.filename, content, category=category
        )
        await db.commit()
        result = {
            "document_id": str(doc.id),
            "filename": doc.filename,
            "space_id": str(doc.space_id),
            "status": doc.status,
            "deduplicated": not is_new,
            "file_hash": hashlib.sha256(content).hexdigest(),
        }
        if is_new:
            from celery_app.tasks import process_document

            task = process_document.apply_async(
                args=(str(doc.id), str(path), str(doc.user_id), str(doc.space_id), doc.category)
            )
            from app.services.rag.knowledge import record_document_enqueue

            await record_document_enqueue(db, str(doc.id), task.id)
            result["celery_task_id"] = task.id
        return result


async def ensure_session(user_id: str, doc_id: str) -> dict:
    """确保办公会话可用：磁盘缓存缺失时从 DB 重建；过期则删除并报错."""
    session = _session_dir(user_id, doc_id)
    if session.exists():
        await _refresh_session_use(user_id, doc_id)
        meta = json.loads(_meta_path(session).read_text(encoding="utf-8"))
        meta["_session"] = str(session)
        return meta
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy import delete, select, update

        async with async_session_factory() as db:
            row = (
                await db.execute(
                    select(OfficeSession).where(
                        OfficeSession.doc_id == str(doc_id),
                        OfficeSession.user_id == uuid.UUID(str(user_id)),
                        OfficeSession.status == "active",
                    )
                )
            ).scalar_one_or_none()
            if row is None:
                raise LookupError("办公文档会话不存在")
            if row.expire_at and row.expire_at < datetime.now(timezone.utc):
                await db.execute(
                    delete(OfficeSession).where(OfficeSession.doc_id == str(doc_id))
                )
                await db.commit()
                raise LookupError("办公文档会话已过期，请重新上传")
            # 重建磁盘缓存（原始文件 + meta）
            ext = Path(row.filename).suffix.lower() or ".txt"
            session.mkdir(parents=True, exist_ok=True)
            if row.file_bytes:
                (session / f"original{ext}").write_bytes(row.file_bytes)
            meta = {
                "doc_id": row.doc_id,
                "filename": row.filename,
                "ext": ext,
                "kind": row.kind,
                "created_at": None,
                "committed": False,
            }
            _meta_path(session).write_text(
                json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
            )
            await db.execute(
                update(OfficeSession)
                .where(OfficeSession.doc_id == str(doc_id))
                .values(last_used_at=datetime.now(timezone.utc), expire_at=_session_expire_at())
            )
            await db.commit()
            meta["_session"] = str(session)
            return meta
    except LookupError:
        raise
    except Exception as exc:  # noqa: BLE001
        logger.warning("办公会话恢复失败: {}", exc)
        raise LookupError("办公文档会话不存在") from exc


async def cleanup_expired_sessions() -> int:
    """清理过期临时会话（DB + 磁盘 + RAG 空间），返回清理数量."""
    removed = 0
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy import select

        async with async_session_factory() as db:
            rows = (
                await db.execute(
                    select(OfficeSession).where(
                        OfficeSession.expire_at.isnot(None),
                        OfficeSession.expire_at < datetime.now(timezone.utc),
                    )
                )
            ).scalars().all()
        for row in rows:
            await discard_session(str(row.user_id), row.doc_id)
            removed += 1
        return removed
    except Exception:  # noqa: BLE001
        return removed


async def discard_session(user_id: str, doc_id: str) -> None:
    """丢弃会话：删除 DB 记录 + RAG 空间数据 + 会话目录."""
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy import delete

        async with async_session_factory() as db:
            await db.execute(
                delete(OfficeSession).where(
                    OfficeSession.doc_id == str(doc_id),
                    OfficeSession.user_id == uuid.UUID(str(user_id)),
                )
            )
            await db.commit()
    except Exception:  # noqa: BLE001
        pass
    session = _session_dir(user_id, doc_id)
    if session.exists():
        try:
            meta = json.loads(_meta_path(session).read_text(encoding="utf-8"))
            if meta.get("space_id"):
                await _delete_space_rows(user_id, meta["space_id"])
        except Exception:  # noqa: BLE001
            pass
    if session.exists():
        shutil.rmtree(session, ignore_errors=True)


def _content_path(meta: dict) -> Path:
    session = Path(meta["_session"])
    buffered = _buffered_path(session, meta["ext"])
    if buffered.exists():
        return buffered
    return _original_path(session, meta["ext"])


def resolve_doc_path(user_id: str, doc_id: str) -> str:
    """解析会话文档的当前文件路径（缓冲优先），供脚本读取."""
    meta = load_session(user_id, doc_id)
    return str(_content_path(meta))


def doc_output_dir(user_id: str, doc_id: str) -> Path:
    """会话输出目录（脚本产物写入这里，可经接口下载）."""
    meta = load_session(user_id, doc_id)
    d = Path(meta["_session"]) / "outputs"
    d.mkdir(parents=True, exist_ok=True)
    return d


def list_doc_outputs(user_id: str, doc_id: str) -> list[dict]:
    """列出会话输出目录里的产物文件（脚本生成的结果，如 csv / 处理后的文件）."""
    try:
        d = doc_output_dir(user_id, doc_id)
        return [
            {"name": f.name, "size": f.stat().st_size}
            for f in sorted(d.iterdir())
            if f.is_file()
        ]
    except Exception:  # noqa: BLE001
        return []


def extract_full_text(user_id: str, doc_id: str) -> str:
    """抽取文档全文（缓冲优先），供 RAG 索引."""
    meta = load_session(user_id, doc_id)
    path = _content_path(meta)
    kind = meta["kind"]
    if kind == "docx":
        return _full_docx(path)
    if kind == "xlsx":
        return _full_xlsx(path)
    if kind == "pptx":
        return _full_pptx(path)
    if kind in {e.lstrip(".") for e in LEGACY_EXTS}:
        return _legacy_extract_text(path, kind)
    if kind in {e.lstrip(".") for e in DOCLING_EXTS}:
        # PDF / 老版 Office / 富文本：走 Docling 统一解析
        from app.services.rag.document_parser import parse_document

        return parse_document(str(path), meta["filename"])
    if kind in {e.lstrip(".") for e in IMAGE_EXTS}:
        # 图片（扫描件/发票/截图）：Docling OCR 提取文本
        from app.services.rag.document_parser import parse_document

        return parse_document(str(path), meta["filename"])
    # 文本类统一走解析入口：CSV/TSV 保留表格列关系，EML 解 MIME 正文，
    # ICS 展开事件字段，同时获得 UTF-8/GB18030/Big5 编码兜底。
    from app.services.rag.document_parser import parse_document

    return parse_document(str(path), meta["filename"])


def _full_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for ti, table in enumerate(doc.tables):
        parts.append(f"[表格{ti}]")
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


def _full_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                parts.append(" | ".join(vals))
    return "\n".join(parts)


def _full_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for si, slide in enumerate(prs.slides):
        parts.append(f"[第{si}页]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)


# ── 结构读取 ─────────────────────────────────────────────

def read_structure(user_id: str, doc_id: str) -> dict:
    """读取文档结构（不是二进制流），供 LLM 决策编辑指令."""
    meta = load_session(user_id, doc_id)
    path = _content_path(meta)
    kind = meta["kind"]
    if kind == "docx":
        structure = _read_docx(path)
    elif kind == "xlsx":
        structure = _read_xlsx(path)
    elif kind == "pptx":
        structure = _read_pptx(path)
    elif kind in {e.lstrip(".") for e in LEGACY_EXTS}:
        # 老版 Office / 富文本：本机 Office COM 提取全文作为结构（只读/分析，不支持结构化编辑）
        structure = _legacy_extract_text(path, kind)[:120000]
    elif kind in {e.lstrip(".") for e in DOCLING_EXTS}:
        # PDF / 老版 Office：Docling 解析全文作为结构（只读/分析，不支持结构化编辑）
        from app.services.rag.document_parser import parse_document

        structure = parse_document(str(path), meta["filename"])[:120000]
    elif kind in {e.lstrip(".") for e in IMAGE_EXTS}:
        # 图片（扫描件/发票/截图）：Docling OCR 提取文本作为结构（只读/分析）
        from app.services.rag.document_parser import parse_document

        structure = parse_document(str(path), meta["filename"])[:120000]
    else:
        # 纯文本（含 EML 邮件 / ICS 日历）：统一解析（编码兜底 + 结构化文本格式化）
        from app.services.rag.document_parser import parse_document

        structure = parse_document(str(path), meta["filename"])[:120000]
    return {"doc_id": doc_id, "kind": kind, "filename": meta["filename"], "structure": structure}


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lines = [f"# 文档段落（共 {len(doc.paragraphs)} 段）"]
    for i, p in enumerate(doc.paragraphs):
        if p.text.strip():
            lines.append(f"[段{i}] {p.text}")
    for ti, table in enumerate(doc.tables):
        lines.append(f"# 表格{ti}（{len(table.rows)}行 x {len(table.columns)}列）")
        for ri, row in enumerate(table.rows[:30]):
            cells = [c.text.strip() for c in row.cells]
            lines.append(f"  行{ri}: {' | '.join(cells)}")
    return "\n".join(lines)[:120000]


def _read_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    lines = [f"# 工作簿（共 {len(wb.sheetnames)} 个表）: {', '.join(wb.sheetnames)}"]
    for ws in wb.worksheets:
        lines.append(f"## Sheet: {ws.title}（{ws.max_row}行 x {ws.max_column}列）")
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 60), values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(vals):
                lines.append("  " + " | ".join(vals))
    return "\n".join(lines)[:120000]


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = [f"# 演示文稿（共 {len(prs.slides)} 页）"]
    for si, slide in enumerate(prs.slides):
        lines.append(f"## 第{si}页")
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    lines.append(f"  [文本] {txt}")
            elif shape.has_table:
                tbl = shape.table
                lines.append(f"  [表格] {len(tbl.rows)}行x{len(tbl.columns)}列")
    return "\n".join(lines)[:120000]


# ── 编辑指令应用 ─────────────────────────────────────────

def apply_edits(user_id: str, doc_id: str, ops: list[dict]) -> list[str]:
    """把编辑指令应用到缓冲副本；返回修改记录列表."""
    meta = load_session(user_id, doc_id)
    session = Path(meta["_session"])
    src = _original_path(session, meta["ext"])
    buffered = _buffered_path(session, meta["ext"])
    if not buffered.exists():
        shutil.copy2(src, buffered)
    records: list[str] = []
    if not (ops or []):
        records.append("⚠️ 未生成任何可执行的编辑操作（文档未修改）")
        return records
    kind = meta["kind"]
    if kind in {e.lstrip(".") for e in IMAGE_EXTS}:
        raise ValueError("图片仅支持读取/分析（OCR 提取文字），不支持结构化编辑")
    for op in ops or []:
        try:
            if kind == "docx":
                r = _apply_docx_op(buffered, op)
            elif kind == "xlsx":
                r = _apply_xlsx_op(buffered, op)
            elif kind == "pptx":
                r = _apply_pptx_op(buffered, op)
            else:
                r = _apply_text_op(buffered, op)
            records.append(r)
        except Exception as exc:  # noqa: BLE001
            records.append(f"❌ 操作 {op.get('op')} 失败：{exc}")
    return records


async def commit_session(user_id: str, doc_id: str) -> dict:
    """把缓冲修改写回真实文件（original），并同步 DB 记录，返回更新后的 meta.

    用户要求"直接利用文件系统修改文档"：编辑成功后立即落盘，
    保存/下载/后续编辑拿到的都是修改后的真实文件。
    """
    meta = load_session(user_id, doc_id)
    session = Path(meta["_session"])
    src = _original_path(session, meta["ext"])
    buffered = _buffered_path(session, meta["ext"])
    if not buffered.exists():
        meta["committed"] = True
        return meta
    # 落盘前备份原始文件（可回滚）
    try:
        backup = session / f"original{meta['ext']}.pre-commit.bak"
        if src.exists() and not backup.exists():
            shutil.copy2(src, backup)
    except Exception:  # noqa: BLE001
        pass
    shutil.copy2(buffered, src)
    try:
        buffered.unlink(missing_ok=True)
    except Exception:  # noqa: BLE001
        pass
    meta["committed"] = True
    _meta_path(session).write_text(
        json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    # 同步 DB：文件字节 + 提取全文（跨实例恢复时保留修改结果）
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import OfficeSession
        from sqlalchemy import update

        new_bytes = src.read_bytes() if src.exists() else None
        new_text = extract_full_text(user_id, doc_id) if new_bytes else None
        async with async_session_factory() as db:
            await db.execute(
                update(OfficeSession)
                .where(
                    OfficeSession.doc_id == str(doc_id),
                    OfficeSession.user_id == uuid.UUID(str(user_id)),
                )
                .values(
                    file_bytes=new_bytes,
                    content_text=new_text,
                    last_used_at=datetime.now(timezone.utc),
                    expire_at=_session_expire_at(),
                )
            )
            await db.commit()
    except Exception as exc:  # noqa: BLE001
        logger.warning("办公会话落盘后 DB 同步失败: {}", exc)
    return meta


def revert_edits(user_id: str, doc_id: str) -> None:
    """撤销未确认的修改：丢弃缓冲副本与落盘前备份，文档保持原样."""
    meta = load_session(user_id, doc_id)
    session = Path(meta["_session"])
    buffered = _buffered_path(session, meta["ext"])
    backup = session / f"original{meta['ext']}.pre-commit.bak"
    for p in (buffered, backup):
        try:
            if p.exists():
                p.unlink()
        except Exception:  # noqa: BLE001
            pass
    meta["committed"] = False
    _meta_path(session).write_text(
        json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def _apply_text_op(path: Path, op: dict) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    op_name = op.get("op")
    if op_name == "rewrite":
        new_text = str(op.get("content") or "")
        path.write_text(new_text, encoding="utf-8")
        return f"✅ 已全文重写（{len(text)}→{len(new_text)} 字符）"
    if op_name == "search_replace":
        old = str(op.get("old") or "")
        new = str(op.get("new") or "")
        if not old:
            raise ValueError("search_replace 缺少 old")
        count = text.count(old)
        if count == 0:
            raise ValueError("未找到要替换的原文")
        path.write_text(text.replace(old, new, 1 if op.get("first_only") else count), encoding="utf-8")
        return f"✅ 已替换 {count} 处：{old[:30]} → {new[:30]}"
    raise ValueError(f"不支持的纯文本操作: {op_name}")


def _apply_docx_op(path: Path, op: dict) -> str:
    from docx import Document

    doc = Document(str(path))
    op_name = op.get("op")
    saved = False

    def _save():
        nonlocal saved
        if not saved:
            doc.save(str(path))
            saved = True

    if op_name == "replace_paragraph":
        idx = int(op.get("index") or 0)
        if idx >= len(doc.paragraphs):
            raise ValueError(f"段落索引越界: {idx}")
        p = doc.paragraphs[idx]
        for run in p.runs:
            run.text = ""
        if p.runs:
            p.runs[0].text = str(op.get("text") or "")
        else:
            p.add_run(str(op.get("text") or ""))
        _save()
        return f"✅ 已替换第{idx}段"
    if op_name == "add_paragraph":
        idx = int(op.get("index") or len(doc.paragraphs))
        p = doc.paragraphs[idx].insert_paragraph_before() if idx < len(doc.paragraphs) else doc.add_paragraph()
        p.add_run(str(op.get("text") or ""))
        _save()
        return "✅ 已新增段落"
    if op_name == "delete_paragraph":
        idx = int(op.get("index") or 0)
        p = doc.paragraphs[idx]._element
        p.getparent().remove(p)
        _save()
        return f"✅ 已删除第{idx}段"
    if op_name == "replace_all":
        find = str(op.get("find") or "")
        replace = str(op.get("replace") or "")
        count = 0
        for p in doc.paragraphs:
            for run in p.runs:
                if find in run.text:
                    run.text = run.text.replace(find, replace)
                    count += 1
        _save()
        return f"✅ 文档正文替换 {find} → {replace}（{count} 处）"
    if op_name == "replace_table_cell":
        ti = int(op.get("table_index") or 0)
        row = int(op.get("row") or 0)
        col = int(op.get("col") or 0)
        cell = doc.tables[ti].cell(row, col)
        cell.text = str(op.get("text") or "")
        _save()
        return f"✅ 表格{ti} 单元格({row},{col}) 已替换"
    if op_name == "add_table_row":
        ti = int(op.get("table_index") or 0)
        values = op.get("values") or []
        row = doc.tables[ti].add_row()
        for ci, v in enumerate(values):
            if ci < len(row.cells):
                row.cells[ci].text = str(v)
        _save()
        return f"✅ 表格{ti} 已新增一行"
    raise ValueError(f"不支持的 docx 操作: {op_name}")


def _apply_xlsx_op(path: Path, op: dict) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path))
    op_name = op.get("op")
    ws = wb[op.get("sheet") or wb.sheetnames[0]] if op.get("sheet") else wb.active
    if op_name == "set_cell":
        ws[op.get("cell")] = op.get("value")
        wb.save(str(path))
        return f"✅ 单元格 {op.get('cell')} = {op.get('value')}"
    if op_name == "replace_all":
        find = str(op.get("find") or "")
        replace = str(op.get("replace") or "")
        count = 0
        for row in ws.iter_rows():
            for c in row:
                if c.value is not None and find in str(c.value):
                    c.value = str(c.value).replace(find, replace)
                    count += 1
        wb.save(str(path))
        return f"✅ Sheet[{ws.title}] 替换 {find} → {replace}（{count} 处）"
    if op_name == "add_row":
        values = op.get("values") or []
        ws.append(values)
        wb.save(str(path))
        return f"✅ Sheet[{ws.title}] 已新增一行"
    raise ValueError(f"不支持的 xlsx 操作: {op_name}")


def _apply_pptx_op(path: Path, op: dict) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    op_name = op.get("op")
    if op_name == "replace_text":
        si = int(op.get("slide_index") or 0)
        find = str(op.get("find") or "")
        replace = str(op.get("replace") or "")
        count = 0
        for shape in prs.slides[si].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
                            count += 1
        prs.save(str(path))
        return f"✅ 第{si}页替换 {find} → {replace}（{count} 处）"
    if op_name == "add_text":
        si = int(op.get("slide_index") or 0)
        slide = prs.slides[si]
        from pptx.util import Inches

        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = str(op.get("text") or "")
        prs.save(str(path))
        return f"✅ 第{si}页已新增文本框"
    if op_name == "delete_slide":
        si = int(op.get("slide_index") or 0)
        xml_slides = prs.slides._sldIdLst
        total = len(xml_slides)
        if si < 0 or si >= total:
            raise ValueError(f"幻灯片索引越界: {si}（共 {total} 页）")
        if total <= 1:
            raise ValueError("不能删除最后一页幻灯片")
        sldId = xml_slides[si]
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            prs.part.drop_rel(rId)
        xml_slides.remove(sldId)
        prs.save(str(path))
        return f"✅ 已删除第{si}页幻灯片（剩余 {total - 1} 页）"
    if op_name == "add_slide":
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题 + 正文版式
        title = str(op.get("title") or "")
        text = str(op.get("text") or "")
        if title and slide.shapes.title is not None:
            slide.shapes.title.text = title
        if text and len(slide.placeholders) > 1:
            slide.placeholders[1].text = text
        prs.save(str(path))
        return f"✅ 已在末尾新增一页幻灯片（共 {len(prs.slides._sldIdLst)} 页）"
    raise ValueError(f"不支持的 pptx 操作: {op_name}")


# ── 会话级 RAG（分析/问答/总结用，复用现有知识库管线） ──────────

def _space_tag(doc_id: str) -> str:
    return f"officedoc_{doc_id}"


async def _delete_space_rows(user_id: str, space_id: str) -> None:
    from sqlalchemy import text

    from app.core.database import async_session_factory

    async with async_session_factory() as session:
        await session.execute(text("DELETE FROM document_chunks WHERE space_id = :i"), {"i": space_id})
        await session.execute(text("DELETE FROM documents WHERE space_id = :i"), {"i": space_id})
        await session.execute(text("DELETE FROM knowledge_spaces WHERE id = :i"), {"i": space_id})
        await session.commit()


async def ensure_rag_index(user_id: str, doc_id: str) -> dict:
    """把文档全文索引进会话专属私有 RAG 空间（幂等，内容变化自动重建）."""
    import hashlib

    from app.core.database import async_session_factory
    from app.services.rag.knowledge import create_space, process_document_pipeline, upload_document_file

    meta = await ensure_session(user_id, doc_id)
    from app.core.executors import run_in_compute

    text = await run_in_compute(extract_full_text, user_id, doc_id)
    if not text.strip():
        raise ValueError("文档无可索引文本")
    content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]
    if meta.get("content_hash") == content_hash and meta.get("space_id"):
        return meta
    if meta.get("space_id"):
        await _delete_space_rows(user_id, meta["space_id"])
    async with async_session_factory() as session:
        space = await create_space(
            session,
            user_id,
            f"办公文档会话-{doc_id}",
            "办公文档分析索引（随会话丢弃）",
            _space_tag(doc_id),
        )
        # Preserve CSV/TSV so the table chunker repeats headers.  Other office
        # formats are already extracted to text/Markdown, so indexing them as
        # Markdown gives the structured chunker headings and tables to work
        # with instead of falling back to a blind .txt split.
        original = Path(str(meta["filename"]))
        index_ext = original.suffix.lower() if original.suffix.lower() in {".csv", ".tsv"} else ".md"
        index_filename = f"{original.stem or 'document'}{index_ext}"
        doc, path, _ = await upload_document_file(
            session,
            user_id,
            str(space.id),
            index_filename,
            text.encode("utf-8"),
        )
        await session.commit()
        await process_document_pipeline(session, str(doc.id), path)
        # The index filename is an implementation detail.  Citations must show
        # the name the user uploaded, not the transient Markdown surrogate.
        doc.filename = str(meta["filename"])
        await session.commit()
    meta["space_id"] = str(space.id)
    meta["rag_doc_id"] = str(doc.id)
    meta["content_hash"] = content_hash
    _meta_path(Path(meta["_session"])).write_text(
        json.dumps(meta, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return meta


async def analyze_doc(
    user_id: str,
    doc_id: str,
    instruction: str,
    mode: str = "qa",
    api_key: str | None = None,
) -> dict:
    """基于会话 RAG 检索 + LLM 作答（总结 / 问答）."""
    from app.agents.skills.base import SkillContext
    from app.core.database import async_session_factory
    from app.services.office_skill_utils import office_llm
    from app.services.rag.query_rewriter import get_retrieval_queries
    from app.services.rag.knowledge import search_user_knowledge

    try:
        await ensure_session(user_id, doc_id)
        space_tag = _space_tag(doc_id)
        office_session = True
    except LookupError:
        # 两条链路兼容：doc_id 也可能是知识空间文档（documents.id），
        # 此时按文档所在知识空间检索，避免"办公文档会话不存在"硬失败。
        office_session = False
        space_tag = await _kb_doc_space_tag(user_id, doc_id)
        if space_tag is None:
            raise LookupError("办公文档会话不存在，且未找到对应知识空间文档") from None

    if office_session:
        from app.core.executors import run_in_compute

        full_text = await run_in_compute(extract_full_text, user_id, doc_id)
        # 小文件全文注入虽避免了检索漏召回，但直接给模型塞到 2 万字符会明显拉高
        # 每个原子步骤的输入 token。上限以模型输入预算为准，超过后走已存在的 RAG 分块。
        direct_limit = min(settings.OFFICE_DOC_FULL_TEXT_LIMIT, 12_000)
        if full_text and len(full_text) <= direct_limit:
            rag_text = full_text
            citations = []
        else:
            await ensure_rag_index(user_id, doc_id)
            query_variants = await get_retrieval_queries(
                instruction, scene="office", user_id=user_id, thinking_mode="think"
            )
            async with async_session_factory() as session:
                rag_text, citations = await search_user_knowledge(
                    session,
                    user_id,
                    query_variants[0] if query_variants else instruction,
                    [space_tag],
                    top_k=6,
                    exclude_categories=["code"],
                    own_space_override=False,
                    rerank_enabled=True,
                    query_variants=query_variants,
                )
    else:
        query_variants = await get_retrieval_queries(
            instruction, scene="office", user_id=user_id, thinking_mode="think"
        )
        async with async_session_factory() as session:
            rag_text, citations = await search_user_knowledge(
                session,
                user_id,
                query_variants[0] if query_variants else instruction,
                [space_tag],
                top_k=6,
                exclude_categories=["code"],
                own_space_override=False,
                rerank_enabled=True,
                query_variants=query_variants,
            )
    if not rag_text:
        return {
            "answer": "文档中未检索到与该问题相关的内容",
            "citations": [],
            "records": [],
        }
    system = (
        "你是文档分析助手。仅依据给定文档片段作答；片段不足以支撑时明确说明，不要编造。"
        if mode == "qa"
        else "你是文档总结助手。基于给定文档片段输出结构化摘要（要点列表，含关键数据）。不要编造。"
    )
    answer = await office_llm(
        SkillContext(user_id=user_id, scene="office", llm_api_key=api_key),
        system,
        f"用户问题/指令：{instruction}\n\n文档片段：\n{rag_text[:24000]}",
        # 文档原子步骤只需提供可供下游汇总的事实，最终整合由后续步骤完成。
        # 限制输出可避免每份文档都生成超长报告，串行 DAG 的等待与 token 成本随之下降。
        max_tokens=1800,
    )
    return {
        "answer": answer,
        "citations": citations,
        "records": [f"已从文档检索到 {len(citations)} 个相关片段"],
    }


async def _kb_doc_space_tag(user_id: str, doc_id: str) -> str | None:
    """知识库文档 → 所在空间 scene_tag（用于按空间检索）；非知识库文档返回 None."""
    try:
        from app.core.database import async_session_factory
        from app.models.db_models import Document, KnowledgeSpace
        from sqlalchemy import select

        async with async_session_factory() as db:
            row = (
                await db.execute(
                    select(KnowledgeSpace.scene_tag)
                    .join(Document, Document.space_id == KnowledgeSpace.id)
                    .where(
                        Document.id == uuid.UUID(str(doc_id)),
                        Document.user_id == uuid.UUID(str(user_id)),
                    )
                )
            ).scalar_one_or_none()
            return row
    except Exception:  # noqa: BLE001
        return None
