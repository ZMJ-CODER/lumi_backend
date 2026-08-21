"""Tests for deterministic new-office-document generation."""

import asyncio

import pytest

from app.services.document_renderer import render_document, safe_output_filename


def test_renderer_creates_docx_pptx_and_xlsx(tmp_path):
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation

    docx = render_document(
        {
            "format": "docx",
            "filename": "方案.docx",
            "title": "项目方案",
            "style": "business",
            "sections": [{"heading": "目标", "paragraphs": ["完成一期交付"], "bullets": ["明确范围"]}],
        },
        tmp_path,
    )
    assert docx.name == "方案.docx"
    assert "项目方案" in "\n".join(p.text for p in Document(str(docx)).paragraphs)

    pptx = render_document(
        {
            "format": "pptx",
            "filename": "汇报.pptx",
            "title": "季度汇报",
            "style": "modern",
            "slides": [{"title": "结论", "bullets": ["收入增长", "成本下降"]}],
        },
        tmp_path,
    )
    presentation = Presentation(str(pptx))
    assert len(presentation.slides) == 1
    assert any("结论" in shape.text for shape in presentation.slides[0].shapes if shape.has_text_frame)

    xlsx = render_document(
        {
            "format": "xlsx",
            "filename": "数据.xlsx",
            "title": "数据表",
            "sheets": [{"name": "汇总", "headers": ["项目", "数值"], "rows": [["收入", "100"]]}],
        },
        tmp_path,
    )
    workbook = load_workbook(str(xlsx), read_only=True)
    assert workbook["汇总"]["A1"].value == "项目"
    assert workbook["汇总"]["B2"].value == "100"


def test_renderer_applies_requested_pptx_style(tmp_path):
    from pptx import Presentation

    path = render_document(
        {
            "format": "pptx",
            "filename": "现代风格.pptx",
            "title": "产品发布",
            "style": "modern",
            "slides": [{"title": "产品发布", "subtitle": "目标与范围"}],
        },
        tmp_path,
    )

    slide = Presentation(str(path)).slides[0]
    assert str(slide.background.fill.fore_color.rgb) == "F0FDFA"
    title_shape = next(shape for shape in slide.shapes if shape.has_text_frame and "产品发布" in shape.text)
    assert str(title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb) == "0F766E"


def test_renderer_rejects_pathlike_filename(tmp_path):
    assert safe_output_filename("方案", "pptx") == "方案.pptx"
    with pytest.raises(ValueError, match="不能包含路径"):
        render_document({"format": "pptx", "filename": "../secret.pptx", "title": "x"}, tmp_path)


def test_document_skill_returns_generic_output_without_path(monkeypatch, tmp_path):
    from app.agents.skills.base import SkillContext
    from app.core.config import settings
    from plugins.skills.office.create_office_document import CreateOfficeDocumentSkill

    monkeypatch.setattr(settings, "UPLOAD_DIR", str(tmp_path / "uploads"))
    result = asyncio.run(
        CreateOfficeDocumentSkill().execute(
            {
                "format": "docx",
                "filename": "交付.docx",
                "title": "交付说明",
                "sections": [{"paragraphs": ["可下载的内容"]}],
            },
            SkillContext(user_id="user-1", conversation_id="job-1", scene="office"),
        )
    )
    assert result.success is True
    output = result.metadata["outputs"][0]
    assert output["name"] == "交付.docx"
    assert output["size"] > 0
    assert output["generic"] is True
    assert "path" not in output
    assert (tmp_path / "uploads" / "office_outputs" / "user-1" / "job-1" / "交付.docx").is_file()


def test_new_document_intent_and_planner_node():
    from app.agents.orchestration.intent import infer_new_office_document, select_named_office_documents
    from app.agents.orchestration.planner import _new_office_document_tree

    document = infer_new_office_document("制作一份现代风格的项目启动演示文稿，生成文件名为启动会.pptx")
    assert document == {"format": "pptx", "filename": "启动会.pptx"}
    assert select_named_office_documents("生成文件名为启动会.pptx", []) == ([], [], False)
    assert infer_new_office_document("将 scores.csv 转为 txt") is None
    tree = _new_office_document_tree("制作演示文稿", document)
    assert tree.nodes[0].agent == "office_document"
    assert tree.nodes[0].params["output_contract"]["expected_output_names"] == ["启动会.pptx"]


def test_document_agent_keeps_planner_file_identity(monkeypatch):
    from app.agents.core.base import WorkerContext
    from app.agents.orchestration.models import TaskNode
    from app.agents.roles.office.agents import OfficeDocumentAgent

    captured = {}

    async def fake_json(*args, **kwargs):
        return {"format": "xlsx", "filename": "wrong.xlsx", "title": "主题", "slides": [{"title": "第一页", "bullets": ["要点"]}]}

    async def fake_run(skill, params, ctx):
        captured["skill"] = skill
        captured["params"] = params
        return {"success": True, "content": "已生成文件：计划.pptx", "outputs": [{"name": "计划.pptx", "size": 1, "generic": True}]}

    monkeypatch.setattr("app.agents.langchain.planning.invoke_json_object", fake_json)
    agent = OfficeDocumentAgent()
    monkeypatch.setattr(agent, "run_skill", fake_run)
    result = asyncio.run(
        agent.execute(
            TaskNode(id="d1", agent="office_document", params={"task": "制作汇报", "format": "pptx", "filename": "计划.pptx"}),
            WorkerContext(user_id="u1", job_id="j1"),
        )
    )
    assert result["success"] is True
    assert captured["skill"] == "create_office_document"
    assert captured["params"]["format"] == "pptx"
    assert captured["params"]["filename"] == "计划.pptx"
