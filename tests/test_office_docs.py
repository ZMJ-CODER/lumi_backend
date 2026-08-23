"""办公文档核心链路测试：格式识别 / 纯文本与 RTF 提取 / 文本与 PPT 编辑 / 撤销."""

import asyncio
from pathlib import Path

import pytest

from app.services import office_docs


def test_document_discovery_and_read_skills_cannot_expand_server_document_scope(monkeypatch):
    from app.agents.skills.base import SkillContext
    from plugins.skills.office.office_docs import InspectDocumentSetSkill, ReadDocumentSkill

    context = SkillContext(user_id="u1", scene="office", office_doc_ids=("allowed",))
    forbidden_read = asyncio.run(ReadDocumentSkill().execute({"doc_id": "other"}, context))
    assert forbidden_read.success is False
    assert forbidden_read.error_code == "FORBIDDEN"

    captured = {}

    async def fake_compute(fn, user_id, doc_ids, query):
        captured["doc_ids"] = doc_ids
        captured["query"] = query
        return [{
            "doc_id": "allowed",
            "filename": "授权文件.pdf",
            "kind": "pdf",
            "summary": "付款条款摘要",
            "page_count": 3,
        }]

    monkeypatch.setattr("plugins.skills.office.office_docs.run_in_compute", fake_compute)
    result = asyncio.run(InspectDocumentSetSkill().execute(
        {"scope": "doc_ids", "doc_ids": ["allowed", "other"], "query": "付款条款"},
        context,
    ))
    assert result.success is True
    assert captured["doc_ids"] == ["allowed"]
    assert result.metadata["document_selection"]["candidate_doc_ids"] == ["allowed"]


def test_detect_kind():
    assert office_docs.detect_kind("a.docx") == "docx"
    assert office_docs.detect_kind("a.xlsx") == "xlsx"
    assert office_docs.detect_kind("a.pptx") == "pptx"
    assert office_docs.detect_kind("a.doc") == "doc"
    assert office_docs.detect_kind("a.pdf") == "pdf"
    assert office_docs.detect_kind("a.txt") == "text"
    assert office_docs.detect_kind("a.md") == "text"


def test_extract_full_text_formats_csv_eml_and_ics(monkeypatch, tmp_path):
    monkeypatch.setattr(office_docs, "OFFICE_DIR", tmp_path / "office")
    user = "u1"

    csv_meta = office_docs.create_session(
        user, "scores.csv", "姓名,成绩\n张三,95\n李四,88".encode("gb18030")
    )
    csv_text = office_docs.extract_full_text(user, csv_meta["doc_id"])
    assert "表头：姓名 | 成绩" in csv_text
    assert "第1行：张三 | 95" in csv_text

    eml = (
        "From: sender@example.com\r\nTo: user@example.com\r\n"
        "Subject: Test mail\r\nContent-Type: text/plain; charset=utf-8\r\n\r\n"
        "邮件正文内容"
    ).encode("utf-8")
    eml_meta = office_docs.create_session(user, "mail.eml", eml)
    eml_text = office_docs.extract_full_text(user, eml_meta["doc_id"])
    assert "发件人：sender@example.com" in eml_text
    assert "邮件正文内容" in eml_text

    ics = (
        "BEGIN:VCALENDAR\r\nBEGIN:VEVENT\r\nSUMMARY:项目会议\r\n"
        "DTSTART;TZID=Asia/Shanghai:20260820T093000\r\n"
        "DESCRIPTION:讨论\\,预算与排期\r\nEND:VEVENT\r\nEND:VCALENDAR\r\n"
    ).encode("utf-8")
    ics_meta = office_docs.create_session(user, "calendar.ics", ics)
    ics_text = office_docs.extract_full_text(user, ics_meta["doc_id"])
    assert "事件1：项目会议" in ics_text
    assert "开始：2026-08-20 09:30" in ics_text
    assert "说明：讨论,预算与排期" in ics_text


def test_extract_doc_text_rtf_masquerade(tmp_path):
    f = tmp_path / "fake.doc"
    f.write_bytes(b"{\\rtf1\\ansi Hello \\b World\\b0 }")
    text = office_docs._extract_doc_text(f)
    assert "Hello" in text


def test_extract_doc_text_plain(tmp_path):
    f = tmp_path / "plain.doc"
    f.write_bytes("纯文本内容 DOCCHAIN1".encode("utf-8"))
    text = office_docs._extract_doc_text(f)
    assert "DOCCHAIN1" in text


def test_apply_edits_text_ops(monkeypatch, tmp_path):
    monkeypatch.setattr(office_docs, "OFFICE_DIR", tmp_path / "office")
    user = "u1"
    meta = office_docs.create_session(user, "a.txt", b"hello world\nsecond line")
    doc_id = meta["doc_id"]
    records = office_docs.apply_edits(
        user, doc_id, [{"op": "search_replace", "old": "world", "new": "lumi", "first_only": True}]
    )
    assert any("已替换" in r for r in records)
    loaded = office_docs.load_session(user, doc_id)
    buffered = office_docs._buffered_path(Path(loaded["_session"]), ".txt")
    assert "hello lumi" in buffered.read_text(encoding="utf-8")

    # 空操作 → 如实报告未修改
    records2 = office_docs.apply_edits(user, doc_id, [])
    assert any("未生成任何可执行的编辑操作" in r for r in records2)


def test_revert_edits_discards_buffer(monkeypatch, tmp_path):
    monkeypatch.setattr(office_docs, "OFFICE_DIR", tmp_path / "office")
    user = "u1"
    meta = office_docs.create_session(user, "a.txt", b"original")
    doc_id = meta["doc_id"]
    office_docs.apply_edits(user, doc_id, [{"op": "rewrite", "content": "changed"}])
    office_docs.revert_edits(user, doc_id)
    meta2 = office_docs.load_session(user, doc_id)
    assert not office_docs._buffered_path(Path(meta2["_session"]), ".txt").exists()


def test_pptx_delete_and_add_slide(tmp_path):
    from pptx import Presentation

    path = tmp_path / "t.pptx"
    prs = Presentation()
    for t in ["一", "二", "三"]:
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = t
    prs.save(str(path))

    r = office_docs._apply_pptx_op(path, {"op": "delete_slide", "slide_index": 1})
    assert "已删除第1页" in r
    prs2 = Presentation(str(path))
    assert [s.shapes.title.text for s in prs2.slides] == ["一", "三"]

    r2 = office_docs._apply_pptx_op(path, {"op": "add_slide", "title": "新页", "text": "内容"})
    assert "已在末尾新增一页" in r2
    prs3 = Presentation(str(path))
    assert len(prs3.slides) == 3
    assert prs3.slides[-1].shapes.title.text == "新页"

    with pytest.raises(ValueError):
        office_docs._apply_pptx_op(path, {"op": "delete_slide", "slide_index": 9})


def test_office_script_agent_passes_doc_ids(monkeypatch):
    from app.agents.core.base import WorkerContext
    from app.agents.orchestration.models import TaskNode
    from app.agents.roles.office.agents import OfficeScriptAgent

    agent = OfficeScriptAgent()
    captured = {}

    async def fake_generate(task, names, ctx, output_contract=None):
        captured["contract"] = output_contract
        return "print('ok')"

    async def fake_run_skill(skill, params, ctx):
        captured["skill"] = skill
        captured["params"] = params
        return {
            "success": True,
            "content": "运行成功",
            "outputs": [{"name": "out.csv", "size": 10, "doc_id": "d1"}],
        }

    monkeypatch.setattr(agent, "_generate_script", fake_generate)
    monkeypatch.setattr(agent, "run_skill", fake_run_skill)
    node = TaskNode(
        id="n1", name="脚本", agent="office_script",
        params={"task": "导出csv", "doc_ids": ["d1", "d2"]}, depends_on=[],
    )
    res = asyncio.run(agent.execute(node, WorkerContext(user_id="u1", job_id="j1")))
    assert res["success"] is True
    assert captured["skill"] == "python_exec"
    assert captured["params"]["doc_ids"] == ["d1", "d2"]
    assert captured["contract"] == {"version": 1, "requires_artifact": False, "expected_output_names": []}
    assert res["outputs"][0]["name"] == "out.csv"


def test_direct_text_conversion_uses_one_selected_document_without_llm(monkeypatch):
    from app.agents.core.base import WorkerContext
    from app.agents.orchestration.models import TaskNode
    from app.agents.roles.office.agents import OfficeScriptAgent

    agent = OfficeScriptAgent()
    captured = {}

    async def forbidden_generate(*args, **kwargs):
        raise AssertionError("简单 CSV 转 TXT 不应调用模型生成脚本")

    async def fake_run_skill(skill, params, ctx):
        captured["skill"] = skill
        captured["params"] = params
        return {
            "success": True,
            "content": "已生成文件：scores.txt",
            "outputs": [{"name": "scores.txt", "size": 42, "doc_id": "scores-doc"}],
        }

    async def fake_ensure(*args, **kwargs):
        return {}

    monkeypatch.setattr(agent, "_generate_script", forbidden_generate)
    monkeypatch.setattr(agent, "run_skill", fake_run_skill)
    monkeypatch.setattr(office_docs, "ensure_session", fake_ensure)
    monkeypatch.setattr(office_docs, "load_session", lambda *_: {"filename": "scores.csv"})
    node = TaskNode(
        id="convert", name="转换", agent="office_script",
        params={
            "task": "将score.csv转为txt",
            "doc_ids": ["scores-doc"],
            "conversion": {
                "source_filename": "scores.csv",
                "target_extension": ".txt",
                "output_filename": "scores.txt",
            },
        },
    )

    result = asyncio.run(agent.execute(node, WorkerContext(user_id="u1", job_id="j1")))

    assert result["success"] is True
    assert result["outputs"] == [{"name": "scores.txt", "size": 42, "doc_id": "scores-doc"}]
    assert captured["skill"] == "python_exec"
    assert captured["params"]["doc_ids"] == ["scores-doc"]
    assert captured["params"]["expected_output_names"] == ["scores.txt"]
    assert "source_name = \"scores.csv\"" in captured["params"]["code"]
    assert "output_name = \"scores.txt\"" in captured["params"]["code"]


def test_direct_text_conversion_script_writes_utf8_output(tmp_path):
    from app.agents.roles.office.agents import OfficeScriptAgent

    source = tmp_path / "scores.csv"
    source.write_bytes("姓名,成绩\n张伟,92\n".encode("gb18030"))
    output_dir = tmp_path / "outputs"
    namespace = {"__name__": "__main__"}
    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setenv("LUMI_DOC_PATHS", '{"scores.csv": "' + str(source).replace("\\", "\\\\") + '"}')
    monkeypatch.setenv("LUMI_DOC_OUTPUT_DIRS", '{"scores.csv": "' + str(output_dir).replace("\\", "\\\\") + '"}')
    try:
        exec(
            OfficeScriptAgent._direct_text_conversion_script(
                {"source_filename": "scores.csv", "target_extension": ".txt", "output_filename": "scores.txt"}
            ),
            namespace,
        )
    finally:
        monkeypatch.undo()
    assert (output_dir / "scores.txt").read_text(encoding="utf-8") == "姓名,成绩\n张伟,92\n"


def test_direct_text_conversion_honors_explicit_tab_delimiter(tmp_path):
    from app.agents.roles.office.agents import OfficeScriptAgent

    source = tmp_path / "scores.csv"
    source.write_text('姓名,备注\n张伟,"语文,补考"\n', encoding="utf-8", newline="")
    output_dir = tmp_path / "outputs"
    namespace = {"__name__": "__main__"}
    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setenv("LUMI_DOC_PATHS", '{"scores.csv": "' + str(source).replace("\\", "\\\\") + '"}')
    monkeypatch.setenv("LUMI_DOC_OUTPUT_DIRS", '{"scores.csv": "' + str(output_dir).replace("\\", "\\\\") + '"}')
    try:
        exec(
            OfficeScriptAgent._direct_text_conversion_script(
                {
                    "source_filename": "scores.csv",
                    "target_extension": ".txt",
                    "output_filename": "scores.txt",
                    "text_delimiter": "\t",
                }
            ),
            namespace,
        )
    finally:
        monkeypatch.undo()
    assert (output_dir / "scores.txt").read_text(encoding="utf-8") == '姓名\t备注\n张伟\t语文,补考\n'


def test_direct_text_conversion_honors_named_output_and_encoding(tmp_path):
    from app.agents.orchestration.intent import resolve_direct_text_conversion
    from app.agents.roles.office.agents import OfficeScriptAgent

    conversion = resolve_direct_text_conversion(
        "将 scores.csv 转为 report.txt，使用 UTF-8-SIG 编码并用逗号分隔",
        [{"doc_id": "scores-doc", "filename": "scores.csv"}],
    )
    assert conversion is not None
    assert conversion["output_filename"] == "report.txt"
    assert conversion["encoding"] == "utf-8-sig"
    assert conversion["text_delimiter"] == ","

    source = tmp_path / "scores.csv"
    source.write_text("姓名,成绩\n张伟,92\n", encoding="utf-8")
    output_dir = tmp_path / "outputs"
    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setenv("LUMI_DOC_PATHS", '{"scores.csv": "' + str(source).replace("\\", "\\\\") + '"}')
    monkeypatch.setenv("LUMI_DOC_OUTPUT_DIRS", '{"scores.csv": "' + str(output_dir).replace("\\", "\\\\") + '"}')
    try:
        exec(OfficeScriptAgent._direct_text_conversion_script(conversion), {"__name__": "__main__"})
    finally:
        monkeypatch.undo()
    assert (output_dir / "report.txt").read_bytes().startswith(b"\xef\xbb\xbf")


def test_python_exec_output_contract_validates_real_text_artifact(tmp_path):
    from plugins.skills.shell.python_exec import _validate_output_contract

    output = tmp_path / "scores.txt"
    output.write_bytes(b"\xef\xbb\xbfname\tscore\nAlice\t92\n")
    contract = {
        "expected_output_names": ["scores.txt"],
        "target_extension": ".txt",
        "encoding": "utf-8-sig",
        "text_delimiter": "\t",
    }
    assert _validate_output_contract(contract, {"scores.txt": [output]}) is None

    output.write_text("name,score\nAlice,92\n", encoding="utf-8")
    assert "UTF-8 BOM" in (_validate_output_contract(contract, {"scores.txt": [output]}) or "")


def test_direct_text_conversion_rejects_success_without_expected_artifact(monkeypatch):
    """脚本 stdout 不能替代真实产物：避免前端展示不存在的下载文件。"""
    from app.agents.core.base import WorkerContext
    from app.agents.orchestration.models import TaskNode
    from app.agents.roles.office.agents import OfficeScriptAgent

    agent = OfficeScriptAgent()

    async def fake_run_skill(*args, **kwargs):
        return {
            "success": False,
            "error": "脚本已结束，但未生成预期文件：scores.txt",
            "error_code": "OUTPUT_MISSING",
            "outputs": [],
        }

    async def fake_ensure(*args, **kwargs):
        return {}

    monkeypatch.setattr(agent, "run_skill", fake_run_skill)
    monkeypatch.setattr(office_docs, "ensure_session", fake_ensure)
    monkeypatch.setattr(office_docs, "load_session", lambda *_: {"filename": "scores.csv"})
    result = asyncio.run(
        agent.execute(
            TaskNode(
                id="convert", agent="office_script",
                params={
                    "task": "将 scores.csv 转为 txt",
                    "doc_ids": ["scores-doc"],
                    "conversion": {"source_filename": "scores.csv", "target_extension": ".txt", "output_filename": "scores.txt"},
                },
            ),
            WorkerContext(user_id="u1", job_id="j1"),
        )
    )

    assert result["success"] is False
    assert result["error_code"] == "OUTPUT_MISSING"


def test_python_exec_marks_sandbox_output_transfer_failure(monkeypatch, tmp_path):
    """回传错误不能退化为通用 EXEC_ERROR，否则 M0 会被错误升级为 M2。"""
    from app.agents.sandbox.base import SandboxResult
    from app.core.config import settings
    from plugins.skills.shell.python_exec import PythonExecSkill

    class FakeSandbox:
        name = "docker"

        async def run_script(self, *args, **kwargs):
            return SandboxResult(status="error", error="沙箱产物回传失败：读取沙箱产物失败")

    monkeypatch.setattr("plugins.skills.shell.python_exec.get_sandbox", lambda: FakeSandbox())
    monkeypatch.setattr(settings, "UPLOAD_DIR", str(tmp_path))
    result = asyncio.run(PythonExecSkill().execute({"code": "print(1)"}))
    assert result.success is False
    assert result.error_code == "SANDBOX_OUTPUT_TRANSFER_FAILED"


def test_analyze_doc_limits_direct_prompt_and_output(monkeypatch):
    captured = {}

    async def fake_ensure(*args, **kwargs):
        return {}

    async def fake_compute(fn, *args):
        return "x" * 12_001

    async def fake_index(*args, **kwargs):
        return {}

    class DummySession:
        async def __aenter__(self):
            return object()

        async def __aexit__(self, *args):
            return False

    async def fake_search(*args, **kwargs):
        return "y" * 30_000, []

    async def fake_llm(context, system, user, **kwargs):
        captured["user"] = user
        captured["max_tokens"] = kwargs["max_tokens"]
        return "answer"

    monkeypatch.setattr(office_docs, "ensure_session", fake_ensure)
    monkeypatch.setattr("app.core.executors.run_in_compute", fake_compute)
    monkeypatch.setattr(office_docs, "ensure_rag_index", fake_index)
    monkeypatch.setattr("app.core.database.async_session_factory", lambda: DummySession())
    monkeypatch.setattr("app.services.rag.knowledge.search_user_knowledge", fake_search)
    monkeypatch.setattr("app.services.office_skill_utils.office_llm", fake_llm)

    result = asyncio.run(office_docs.analyze_doc("u1", "d1", "总结"))

    assert result["answer"] == "answer"
    assert captured["max_tokens"] == 1800
    assert len(captured["user"].split("文档片段：\n", 1)[1]) == 24_000


def test_office_script_generation_uses_one_llm_call(monkeypatch):
    from app.agents.core.base import WorkerContext
    from app.agents.roles.office.agents import OfficeScriptAgent

    calls = []

    class FakeLLM:
        async def chat(self, messages, **kwargs):
            calls.append((messages, kwargs))
            return "print('ok')"

    monkeypatch.setattr("app.core.llm.LLMClient", FakeLLM)

    code = asyncio.run(
        OfficeScriptAgent()._generate_script(
            "把成绩表导出为摘要 CSV", ["scores.csv"], WorkerContext(user_id="u1", job_id="j1")
        )
    )

    assert code == "print('ok')"
    assert len(calls) == 1
    assert calls[0][1]["max_tokens"] == 4000
    assert "输出契约" in calls[0][0][0]["content"]


def test_office_script_generation_retries_once_after_empty_model_reply(monkeypatch):
    from app.agents.core.base import WorkerContext
    from app.agents.roles.office.agents import OfficeScriptAgent

    calls = 0

    class FakeLLM:
        async def chat(self, messages, **kwargs):
            nonlocal calls
            calls += 1
            if calls == 1:
                raise RuntimeError("模型返回空内容")
            return "print('ok')"

    monkeypatch.setattr("app.core.llm.LLMClient", FakeLLM)
    code = asyncio.run(
        OfficeScriptAgent()._generate_script("导出 CSV", ["scores.csv"], WorkerContext(user_id="u1", job_id="j1"))
    )
    assert code == "print('ok')"
    assert calls == 2
